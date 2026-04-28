import streamlit as st
import requests
import json
from io import BytesIO
import pandas as pd
import numpy as np
import plotly.express as px

# -----------------------------
# SAFE IMPORTS
# -----------------------------
try:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from pptx import Presentation
    LIBS_AVAILABLE = True
except Exception:
    LIBS_AVAILABLE = False


# -----------------------------
# CONFIG
# -----------------------------
API_URL = "https://emea.snaplogic.com/api/1/rest/slsched/feed/ConnectFasterInc/projects/Tejasri%20Reddy%20Beeram/Diabetes_v5%20Task"
API_KEY = "681vyxaddAuuuJ21FD7LOYVogX7B0dHB"

st.set_page_config("Pharma Intelligence Hub", "💊", layout="wide")


# -----------------------------
# PDF
# -----------------------------
def create_pdf(text):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()

    story = []
    for line in text.split("\n"):
        if line.strip():
            story.append(Paragraph(line, styles["Normal"]))
            story.append(Spacer(1, 6))

    doc.build(story)
    buffer.seek(0)
    return buffer


# -----------------------------
# PPT
# -----------------------------
def create_ppt(text):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Pharma Intelligence Report"
    slide.placeholders[1].text = "Executive Dashboard Output"

    for sec in text.split("\n")[:8]:
        if sec.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Insight"
            slide.placeholders[1].text = sec[:900]

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# -----------------------------
# REGION PARSER (CORE ENGINE)
# -----------------------------
def parse_regions(text):

    mapping = {
        "Riyadh": 4,
        "Jeddah": 3,
        "Dammam": 3,
        "Makkah": 2,
        "Madinah": 2,
        "Buraidah": 2,
        "Khobar": 2,
        "Other": 1
    }

    data = []

    for region, score in mapping.items():
        if region.lower() in text.lower():
            data.append({
                "Region": region,
                "Score": score
            })

    return pd.DataFrame(data) if data else None


# -----------------------------
# BRAND SEGMENTATION (SIMULATED)
# -----------------------------
def brand_segmentation():
    return pd.DataFrame({
        "Brand": ["Empagliflozin", "Metformin", "Insulin", "GLP-1"],
        "Priority Score": [4, 3, 3, 2],
        "Market Share Impact": [35, 25, 30, 10]
    })


# -----------------------------
# UI
# -----------------------------
st.title("💊 Pharma Intelligence Hub (KSA)")

with st.form("form"):
    prompt = st.text_area("Enter Query", height=120)

    role = st.radio(
        "Role",
        ["Sales Representative", "Market Access Specialist", "Medical Science Liaison (MSL)"]
    )

    submit = st.form_submit_button("Generate Intelligence")


# -----------------------------
# API CALL
# -----------------------------
if submit and prompt.strip():

    with st.spinner("Generating Pharma Intelligence..."):

        payload = [{"question_type": role, "prompt": prompt}]
        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }

        res = requests.post(API_URL, headers=headers, data=json.dumps(payload))
        data = res.json()

        reply = data[0]["Customer_Story"] if isinstance(data, list) else str(data)


    # -----------------------------
    # RESPONSE
    # -----------------------------
    st.subheader("🧠 AI Intelligence Output")
    st.info(reply)


    # =========================================================
    # 📊 KPI DASHBOARD (POWER BI STYLE)
    # =========================================================
    st.markdown("## 📊 Executive KPI Dashboard")

    df_region = parse_regions(reply)

    col1, col2, col3 = st.columns(3)

    if df_region is not None:
        with col1:
            st.metric("High Priority Hubs", len(df_region[df_region["Score"] >= 3]))

        with col2:
            st.metric("Avg Priority Score", round(df_region["Score"].mean(), 2))

        with col3:
            st.metric("Total Regions", len(df_region))


    # =========================================================
    # 🗺️ SAUDI MAP (HEATMAP STYLE)
    # =========================================================
    st.markdown("## 🗺️ Saudi Arabia Priority Heatmap")

    if df_region is not None:

        # Fake coordinates (industry-standard approach)
        coords = {
            "Riyadh": [24.7136, 46.6753],
            "Jeddah": [21.4858, 39.1925],
            "Dammam": [26.4207, 50.0888],
            "Makkah": [21.3891, 39.8579],
            "Madinah": [24.5247, 39.5692],
            "Buraidah": [26.3260, 43.9750],
            "Khobar": [26.2172, 50.1971],
        }

        df_region["lat"] = df_region["Region"].map(lambda x: coords.get(x, [24, 45])[0])
        df_region["lon"] = df_region["Region"].map(lambda x: coords.get(x, [24, 45])[1])

        fig = px.scatter_mapbox(
            df_region,
            lat="lat",
            lon="lon",
            size="Score",
            color="Score",
            hover_name="Region",
            zoom=4,
            mapbox_style="open-street-map"
        )

        st.plotly_chart(fig, use_container_width=True)


    # =========================================================
    # 🧬 BRAND SEGMENTATION
    # =========================================================
    st.markdown("## 🧬 Brand Performance Segmentation")

    df_brand = brand_segmentation()

    col1, col2 = st.columns(2)

    with col1:
        st.bar_chart(df_brand.set_index("Brand")["Priority Score"])

    with col2:
        st.bar_chart(df_brand.set_index("Brand")["Market Share Impact"])


    # =========================================================
    # DOWNLOADS
    # =========================================================
    st.markdown("## ⬇️ Reports")

    if LIBS_AVAILABLE:
        pdf = create_pdf(reply)
        ppt = create_ppt(reply)

        c1, c2 = st.columns(2)

        with c1:
            st.download_button("📄 PDF", pdf, "Pharma.pdf", "application/pdf")

        with c2:
            st.download_button("📊 PPT", ppt, "Pharma.pptx",
                               "application/vnd.openxmlformats-officedocument.presentationml.presentation")

else:
    st.caption("Enter query to generate pharma intelligence")
